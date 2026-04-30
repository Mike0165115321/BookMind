import os
import json
import fitz  # PyMuPDF
import docx  # python-docx
import xml.etree.ElementTree as ET
import pandas as pd
from pptx import Presentation
from bs4 import BeautifulSoup

class DocumentLoader:
    """
    Dedicated class to handle text extraction from various file formats.
    """
    @staticmethod
    def load(filepath, book_title=None):
        filename = os.path.basename(filepath).lower()
        
        if filename.endswith((".jsonl", ".json")):
            return DocumentLoader._load_json(filepath, book_title)
        elif filename.endswith(".pdf"):
            return DocumentLoader._load_pdf(filepath, book_title)
        elif filename.endswith((".docx", ".doc")):
            return DocumentLoader._load_docx(filepath, book_title)
        elif filename.endswith(".svg"):
            return DocumentLoader._load_svg(filepath, book_title)
        elif filename.endswith((".xlsx", ".xls", ".csv")):
            return DocumentLoader._load_excel_csv(filepath, book_title)
        elif filename.endswith(".pptx"):
            return DocumentLoader._load_pptx(filepath, book_title)
        elif filename.endswith((".html", ".htm")):
            return DocumentLoader._load_html(filepath, book_title)
        elif filename.endswith((".md", ".markdown", ".txt")):
            return DocumentLoader._load_text(filepath, book_title)
        else:
            return DocumentLoader._load_text(filepath, book_title)

    @staticmethod
    def _load_json(filepath, book_title=None):
        docs = []
        try:
            with open(filepath, "r", encoding="utf-8") as f:
                content = f.read().strip()
                if content.startswith('[') and content.endswith(']'):
                    data = json.loads(content)
                    for obj in data:
                        docs.append(DocumentLoader._format_json_obj(obj, book_title))
                    return docs
        except:
            pass

        with open(filepath, "r", encoding="utf-8") as f:
            for line in f:
                line = line.strip()
                if not line: continue
                try:
                    obj = json.loads(line)
                    docs.append(DocumentLoader._format_json_obj(obj, book_title))
                except: continue
        return docs

    @staticmethod
    def _format_json_obj(obj, book_title=None):
        b_title = book_title or obj.get("book_title", "")
        title = obj.get("title", "")
        content = obj.get("content", "")
        
        prefix = ""
        if b_title: prefix += f"[{b_title}] "
        if title: prefix += f"{title}\n"
        
        return {"content": content, "metadata_prefix": prefix.strip()}

    @staticmethod
    def _load_pdf(filepath, book_title=None):
        text = ""
        try:
            with fitz.open(filepath) as doc:
                for page in doc:
                    text += page.get_text() + "\n"
        except Exception as e:
            print(f"❌ Error loading PDF: {e}")
        return [{"content": text, "metadata_prefix": f"[{book_title}]" if book_title else ""}]

    @staticmethod
    def _load_docx(filepath, book_title=None):
        text = ""
        try:
            doc = docx.Document(filepath)
            for para in doc.paragraphs:
                text += para.text + "\n"
        except Exception as e:
            print(f"❌ Error loading DOCX: {e}")
        return [{"content": text, "metadata_prefix": f"[{book_title}]" if book_title else ""}]

    @staticmethod
    def _load_svg(filepath, book_title=None):
        text_elements = []
        try:
            tree = ET.parse(filepath)
            root = tree.getroot()
            # Namespace handling for SVG
            ns = {'svg': 'http://www.w3.org/2000/svg'}
            # Find all text-related elements
            for elem in root.findall('.//svg:text', ns):
                if elem.text:
                    text_elements.append(elem.text.strip())
                for tspan in elem.findall('.//svg:tspan', ns):
                    if tspan.text:
                        text_elements.append(tspan.text.strip())
            
            combined_text = "\n".join(text_elements)
        except Exception as e:
            print(f"❌ Error loading SVG: {e}")
            combined_text = ""
        
        return [{"content": combined_text, "metadata_prefix": f"[{book_title}]" if book_title else ""}]

    @staticmethod
    def _load_excel_csv(filepath, book_title=None):
        text = ""
        try:
            if filepath.endswith(".csv"):
                df = pd.read_csv(filepath)
            else:
                df = pd.read_excel(filepath)
            
            # Convert to a descriptive string for RAG
            text = f"ข้อมูลตารางจากหนังสือ: {book_title if book_title else 'เอกสารข้อมูล'}\n"
            text += df.to_string(index=False)
        except Exception as e:
            print(f"❌ Error loading Excel/CSV: {e}")
        
        return [{"content": text, "metadata_prefix": f"[{book_title}]" if book_title else ""}]

    @staticmethod
    def _load_pptx(filepath, book_title=None):
        text = ""
        try:
            prs = Presentation(filepath)
            for i, slide in enumerate(prs.slides):
                text += f"\n--- สไลด์ที่ {i+1} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        except Exception as e:
            print(f"❌ Error loading PPTX: {e}")
        return [{"content": text, "metadata_prefix": f"[{book_title}]" if book_title else ""}]

    @staticmethod
    def _load_html(filepath, book_title=None):
        text = ""
        try:
            with open(filepath, "r", encoding="utf-8") as f:
                soup = BeautifulSoup(f.read(), "html.parser")
                # Remove script and style elements
                for script_or_style in soup(["script", "style"]):
                    script_or_style.decompose()
                text = soup.get_text(separator="\n")
        except Exception as e:
            print(f"❌ Error loading HTML: {e}")
        return [{"content": text, "metadata_prefix": f"[{book_title}]" if book_title else ""}]

    @staticmethod
    def _load_text(filepath, book_title=None):
        try:
            with open(filepath, "r", encoding="utf-8") as f:
                return [{"content": f.read(), "metadata_prefix": f"[{book_title}]" if book_title else ""}]
        except:
            return []
